"""
Script pour générer une présentation PowerPoint pour Mazelet Blasa
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RgbColor
import os

# Chemins des images
IMAGES_DIR = r"C:\Users\pc\.gemini\antigravity\brain\d8e85880-abdf-461a-a250-1a8af2114102"
ARCHITECTURE_IMG = os.path.join(IMAGES_DIR, "uploaded_image_0_1765040011707.jpg")
UML_IMG = os.path.join(IMAGES_DIR, "uploaded_image_1_1765040011707.png")
ER_IMG = os.path.join(IMAGES_DIR, "uploaded_image_2_1765040011707.png")

def set_slide_background(slide, r, g, b):
    """Définir la couleur de fond d'une slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RgbColor(r, g, b)

def add_title_slide(prs, title, subtitle):
    """Créer une slide de titre"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 0, 102, 153)  # Bleu foncé
    
    # Titre principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Sous-titre
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, image_path=None):
    """Créer une slide avec contenu"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 255, 255, 255)
    
    # Barre de titre
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RgbColor(0, 102, 153)
    title_shape.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Contenu
    if image_path and os.path.exists(image_path):
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
        slide.shapes.add_picture(image_path, Inches(4.8), Inches(1.5), width=Inches(4.8))
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(51, 51, 51)
        p.space_after = Pt(12)
    
    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """Créer une slide centrée sur une image"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 255, 255, 255)
    
    # Barre de titre
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RgbColor(0, 102, 153)
    title_shape.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Image centrée
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.2), width=Inches(9))
    
    # Caption
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.5))
        tf = caption_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = RgbColor(51, 51, 51)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_features_slide(prs):
    """Slide des fonctionnalités avec icônes"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 255, 255, 255)
    
    # Barre de titre
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RgbColor(0, 102, 153)
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Fonctionnalites Principales"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    features = [
        ("Gestion des matchs", "Creation, modification, suppression"),
        ("Authentification JWT", "Login/Signup securise"),
        ("Participation", "Rejoindre/Quitter des matchs"),
        ("Filtres avances", "Ville, Prix, Date"),
        ("Recommandations KNN", "Suggestions personnalisees"),
        ("Alertes Email", "Rappels et annulations"),
    ]
    
    for i, (ftitle, desc) in enumerate(features):
        row = i // 2
        col = i % 2
        x = Inches(0.5 + col * 4.8)
        y = Inches(1.3 + row * 1.7)
        
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.5), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(240, 248, 255)
        box.line.color.rgb = RgbColor(0, 102, 153)
        
        # Titre feature
        feat_title = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(4), Inches(0.5))
        tf = feat_title.text_frame
        p = tf.paragraphs[0]
        p.text = ftitle
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RgbColor(0, 102, 153)
        
        # Description
        feat_desc = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), Inches(4), Inches(0.6))
        tf = feat_desc.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(51, 51, 51)
    
    return slide

def add_conclusion_slide(prs):
    """Slide de conclusion"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 0, 102, 153)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Conclusion"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Points clés
    points = [
        "Facilite l'organisation de matchs de football",
        "Experience personnalisee grace a l'IA (KNN)",
        "Stack technique moderne (React + FastAPI + PostgreSQL)",
        "Securite avec JWT et alertes email"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(3))
    tf = content_box.text_frame
    
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "-> " + point
        p.font.size = Pt(22)
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.space_after = Pt(20)
    
    # Remerciements
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.8))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Merci de votre attention !"
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Créer la présentation complète"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Titre
    add_title_slide(prs, "Mazelet Blasa", 
                   "Plateforme de gestion et recommandation de matchs de football\nSysteme intelligent base sur KNN")
    
    # Slide 2: Problème
    add_content_slide(prs, "Le Probleme", [
        "Difficulte a trouver des matchs adaptes a ses preferences",
        "Navigation inefficace parmi de nombreux matchs",
        "Manque de personnalisation dans l'experience utilisateur",
        "Pas de suggestions basees sur l'historique des participations",
        "Temps perdu en recherche manuelle"
    ])
    
    # Slide 3: Solution
    add_content_slide(prs, "Notre Solution", [
        "Plateforme web moderne et intuitive",
        "Systeme de recommandation intelligent (KNN)",
        "Authentification securisee avec JWT",
        "Alertes email automatiques (rappels, annulations)",
        "Interface responsive et elegante"
    ])
    
    # Slide 4: Architecture
    add_image_slide(prs, "Architecture du Systeme", ARCHITECTURE_IMG,
                   "Frontend React + Backend FastAPI + Base de donnees PostgreSQL")
    
    # Slide 5: Diagramme de classes UML
    add_image_slide(prs, "Diagramme de Classes UML", UML_IMG,
                   "Modelisation orientee objet du systeme")
    
    # Slide 6: Modèle Entité-Relation
    add_image_slide(prs, "Modele Entite-Relation (MCD)", ER_IMG,
                   "Structure de la base de donnees relationnelle")
    
    # Slide 7: Algorithme KNN
    add_content_slide(prs, "Algorithme KNN - Recommandations", [
        "Analyse l'historique de participation de l'utilisateur",
        "Encode les caracteristiques: Ville, Stade, Type, Nb joueurs",
        "Calcule la distance euclidienne entre matchs",
        "Recommande les matchs les plus similaires",
        "Affiche un score de similarite (ex: 85% match)"
    ])
    
    # Slide 8: Fonctionnalités
    add_features_slide(prs)
    
    # Slide 9: Technologies
    add_content_slide(prs, "Stack Technique", [
        "Frontend: React, JavaScript, CSS moderne",
        "Backend: FastAPI (Python), Pydantic, JWT",
        "Base de donnees: PostgreSQL, SQLAlchemy ORM",
        "Machine Learning: scikit-learn, NumPy",
        "Services: SMTP Gmail pour les alertes email"
    ])
    
    # Slide 10: Perspectives
    add_content_slide(prs, "Perspectives d'Amelioration", [
        "Systeme de notation des matchs",
        "Chat/Messagerie entre joueurs",
        "Statistiques avancees pour les profils",
        "Integration de paiement en ligne",
        "Notifications push (Mobile App)"
    ])
    
    # Slide 11: Conclusion
    add_conclusion_slide(prs)
    
    # Sauvegarder
    output_path = r"c:\Users\pc\Downloads\learn react - Copie\Mazelet_Blasa_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation creee avec succes: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
