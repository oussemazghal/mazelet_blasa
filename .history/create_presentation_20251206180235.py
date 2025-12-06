"""
Script pour generer une presentation PowerPoint pour Mazelet Blasa
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Chemins des images
IMAGES_DIR = r"C:\Users\pc\.gemini\antigravity\brain\d8e85880-abdf-461a-a250-1a8af2114102"
ARCHITECTURE_IMG = os.path.join(IMAGES_DIR, "uploaded_image_0_1765040011707.jpg")
UML_IMG = os.path.join(IMAGES_DIR, "uploaded_image_1_1765040011707.png")
ER_IMG = os.path.join(IMAGES_DIR, "uploaded_image_2_1765040011707.png")

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def set_font_color(font, hex_color):
    """Set font color using hex"""
    from pptx.dml.color import RgbColor
    r, g, b = hex_to_rgb(hex_color)
    font.color.rgb = RgbColor(r, g, b)

def set_fill_color(fill, hex_color):
    """Set fill color using hex"""
    from pptx.dml.color import RgbColor
    r, g, b = hex_to_rgb(hex_color)
    fill.solid()
    fill.fore_color.rgb = RgbColor(r, g, b)

# Couleurs
PRIMARY = "#006699"
WHITE = "#FFFFFF"
DARK = "#333333"
LIGHT_BLUE = "#F0F8FF"

def add_title_slide(prs, title, subtitle):
    """Creer une slide de titre"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Fond
    set_fill_color(slide.background.fill, PRIMARY)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    set_font_color(p.font, WHITE)
    p.alignment = PP_ALIGN.CENTER
    
    # Sous-titre
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    set_font_color(p.font, WHITE)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, items):
    """Slide avec contenu"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_fill_color(slide.background.fill, WHITE)
    
    # Barre titre
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    set_fill_color(bar.fill, PRIMARY)
    bar.line.fill.background()
    
    # Titre
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    set_font_color(p.font, WHITE)
    
    # Contenu
    cb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = cb.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "-> " + item
        p.font.size = Pt(18)
        set_font_color(p.font, DARK)
        p.space_after = Pt(12)
    
    return slide

def add_image_slide(prs, title, img_path, caption=""):
    """Slide avec image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_fill_color(slide.background.fill, WHITE)
    
    # Barre titre
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    set_fill_color(bar.fill, PRIMARY)
    bar.line.fill.background()
    
    # Titre
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    set_font_color(p.font, WHITE)
    
    # Image
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.2), width=Inches(9))
    
    # Caption
    if caption:
        cb = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.5))
        p = cb.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        set_font_color(p.font, DARK)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_features_slide(prs):
    """Slide fonctionnalites"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_fill_color(slide.background.fill, WHITE)
    
    # Barre titre
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    set_fill_color(bar.fill, PRIMARY)
    bar.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = "Fonctionnalites Principales"
    p.font.size = Pt(28)
    p.font.bold = True
    set_font_color(p.font, WHITE)
    
    features = [
        ("Gestion des matchs", "Creation, modification, suppression"),
        ("Authentification JWT", "Login/Signup securise"),
        ("Participation", "Rejoindre/Quitter des matchs"),
        ("Filtres avances", "Ville, Prix, Date"),
        ("Recommandations KNN", "Suggestions personnalisees"),
        ("Alertes Email", "Rappels et annulations"),
    ]
    
    for i, (ftitle, desc) in enumerate(features):
        row, col = i // 2, i % 2
        x, y = Inches(0.5 + col * 4.8), Inches(1.3 + row * 1.7)
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.5), Inches(1.5))
        set_fill_color(box.fill, LIGHT_BLUE)
        
        ft = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(4), Inches(0.5))
        p = ft.text_frame.paragraphs[0]
        p.text = ftitle
        p.font.size = Pt(18)
        p.font.bold = True
        set_font_color(p.font, PRIMARY)
        
        fd = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), Inches(4), Inches(0.6))
        p = fd.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        set_font_color(p.font, DARK)
    
    return slide

def add_conclusion_slide(prs):
    """Slide conclusion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_fill_color(slide.background.fill, PRIMARY)
    
    # Titre
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Conclusion"
    p.font.size = Pt(40)
    p.font.bold = True
    set_font_color(p.font, WHITE)
    p.alignment = PP_ALIGN.CENTER
    
    # Points
    points = [
        "Facilite l'organisation de matchs de football",
        "Experience personnalisee grace a KNN",
        "Stack moderne: React + FastAPI + PostgreSQL",
        "Securite avec JWT et alertes email"
    ]
    
    cb = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(3))
    tf = cb.text_frame
    
    for i, pt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "-> " + pt
        p.font.size = Pt(22)
        set_font_color(p.font, WHITE)
        p.space_after = Pt(20)
    
    # Merci
    thx = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.8))
    p = thx.text_frame.paragraphs[0]
    p.text = "Merci de votre attention!"
    p.font.size = Pt(28)
    p.font.italic = True
    set_font_color(p.font, WHITE)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slides
    add_title_slide(prs, "MAZELET BLASA", 
        "Plateforme de gestion de matchs de football\nSysteme de recommandation KNN")
    
    add_content_slide(prs, "Le Probleme", [
        "Difficulte a trouver des matchs adaptes",
        "Navigation inefficace parmi de nombreux matchs",
        "Manque de personnalisation",
        "Pas de suggestions basees sur l'historique",
        "Temps perdu en recherche manuelle"
    ])
    
    add_content_slide(prs, "Notre Solution", [
        "Plateforme web moderne et intuitive",
        "Systeme de recommandation KNN",
        "Authentification securisee JWT",
        "Alertes email automatiques",
        "Interface responsive"
    ])
    
    add_image_slide(prs, "Architecture du Systeme", ARCHITECTURE_IMG,
        "React + FastAPI + PostgreSQL")
    
    add_image_slide(prs, "Diagramme de Classes UML", UML_IMG,
        "Modelisation orientee objet")
    
    add_image_slide(prs, "Modele Entite-Relation", ER_IMG,
        "Structure de la base de donnees")
    
    add_content_slide(prs, "Algorithme KNN", [
        "Analyse l'historique de participation",
        "Encode: Ville, Stade, Type, Nb joueurs",
        "Calcule la distance euclidienne",
        "Recommande les matchs similaires",
        "Affiche un score de similarite"
    ])
    
    add_features_slide(prs)
    
    add_content_slide(prs, "Stack Technique", [
        "Frontend: React, JavaScript, CSS",
        "Backend: FastAPI, Pydantic, JWT",
        "Database: PostgreSQL, SQLAlchemy",
        "ML: scikit-learn, NumPy",
        "Email: SMTP Gmail"
    ])
    
    add_content_slide(prs, "Perspectives", [
        "Systeme de notation des matchs",
        "Chat entre joueurs",
        "Statistiques avancees",
        "Paiement en ligne",
        "Notifications push"
    ])
    
    add_conclusion_slide(prs)
    
    # Save
    out = r"c:\Users\pc\Downloads\learn react - Copie\Mazelet_Blasa_Presentation.pptx"
    prs.save(out)
    print(f"Presentation creee: {out}")

if __name__ == "__main__":
    main()
